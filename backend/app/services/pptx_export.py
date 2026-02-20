"""PPTX export service — generates a 16:9 PowerPoint from DeckOut."""

import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

from app.schemas.common import BulletOut, DeckOut

# Slide dimensions: 16:9
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Color palette ──
PRIMARY = RGBColor(0x1A, 0x56, 0xDB)      # primary blue
PRIMARY_DARK = RGBColor(0x12, 0x3E, 0xA8)  # darker blue for gradients
ACCENT = RGBColor(0x00, 0x96, 0x88)        # teal accent
DARK = RGBColor(0x1E, 0x1E, 0x2E)          # near-black for body text
HEADING = RGBColor(0x11, 0x11, 0x27)       # slightly darker for headings
GRAY = RGBColor(0x6C, 0x6C, 0x80)          # secondary text
LIGHT_GRAY = RGBColor(0x9E, 0x9E, 0xAE)   # tertiary / muted text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xFA)     # subtle background
FORMULA_COLOR = RGBColor(0x7B, 0x1F, 0xA2) # purple for formula bullets
DIVIDER_LINE = RGBColor(0xE0, 0xE0, 0xEB)  # thin rule color

# Bullet prefixes by content type
BULLET_CHAR = "\u2022"       # default: •
FORMULA_CHAR = "\U0001D453"  # 𝑓 (mathematical italic f)

MAX_BULLETS_PER_SLIDE = 6


# ── Helpers ──

def _add_textbox(slide, left, top, width, height, text,
                 font_size=18, color=DARK, bold=False, italic=False,
                 alignment=None, font_name="Calibri"):
    """Add a textbox with a single styled paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    if alignment is not None:
        p.alignment = alignment
    return txBox


def _add_filled_rect(slide, left, top, width, height, fill_color):
    """Add a solid-filled rectangle shape (no outline)."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_gradient_rect(slide, left, top, width, height, color1, color2):
    """Add a rectangle with a two-stop linear gradient fill."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height,
    )
    shape.line.fill.background()

    # Build gradient XML
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 0  # left to right

    stops = fill._fill.xpath('.//a:gsLst/a:gs', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
    if len(stops) >= 2:
        stops[0].find(qn('a:srgbClr')).set('val', f'{color1.red:02X}{color1.green:02X}{color1.blue:02X}' if hasattr(color1, 'red') else str(color1))
        stops[1].find(qn('a:srgbClr')).set('val', f'{color2.red:02X}{color2.green:02X}{color2.blue:02X}' if hasattr(color2, 'red') else str(color2))

    return shape


def _add_thin_rule(slide, left, top, width, color=DIVIDER_LINE):
    """Add a thin horizontal line."""
    shape = slide.shapes.add_shape(
        1,  # rectangle used as line
        left, top, width, Pt(1.5),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_slide_number(slide, slide_num: int, total: int):
    """Add a small slide number indicator at bottom-right."""
    _add_textbox(
        slide,
        Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.35),
        f"{slide_num} / {total}",
        font_size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT,
    )


def _add_section_badge(slide, section_heading: str, section_idx: int, total_sections: int):
    """Add a subtle section context badge at top-left with progress dots."""
    # Section label
    _add_textbox(
        slide,
        Inches(0.8), Inches(0.35), Inches(5), Inches(0.35),
        section_heading,
        font_size=11, color=LIGHT_GRAY, font_name="Calibri",
    )
    # Progress dots — filled for current and past sections, hollow for future
    dots = ""
    for i in range(total_sections):
        dots += "\u25CF " if i <= section_idx else "\u25CB "
    _add_textbox(
        slide,
        Inches(0.8), Inches(0.6), Inches(5), Inches(0.25),
        dots.strip(),
        font_size=7, color=LIGHT_GRAY,
    )


# ── Slide builders ──

def _add_title_slide(prs, title: str):
    """Title page with accent stripe."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Left accent stripe
    _add_filled_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, PRIMARY)

    # Title
    _add_textbox(
        slide,
        Inches(1.2), Inches(2.2), Inches(11), Inches(2.0),
        title, font_size=40, color=HEADING, bold=True,
    )

    # Subtitle
    _add_textbox(
        slide,
        Inches(1.2), Inches(4.5), Inches(11), Inches(0.6),
        "Generated by SlideNode", font_size=16, color=LIGHT_GRAY, italic=True,
    )

    # Bottom rule
    _add_thin_rule(slide, Inches(1.2), Inches(6.8), Inches(4))


def _add_section_slide(prs, heading: str, summary_note: str, section_idx: int, total_sections: int):
    """Section divider — large heading with accent bar and progress."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top accent bar
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), PRIMARY)

    # Section number
    section_label = f"SECTION {section_idx + 1} OF {total_sections}"
    _add_textbox(
        slide,
        Inches(1), Inches(1.8), Inches(11.333), Inches(0.5),
        section_label, font_size=13, color=ACCENT, bold=True,
    )

    # Heading
    _add_textbox(
        slide,
        Inches(1), Inches(2.5), Inches(11.333), Inches(1.5),
        heading, font_size=36, color=HEADING, bold=True,
    )

    # Divider line
    _add_thin_rule(slide, Inches(1), Inches(4.2), Inches(3), color=PRIMARY)

    # Summary note
    if summary_note:
        _add_textbox(
            slide,
            Inches(1), Inches(4.5), Inches(11.333), Inches(1.2),
            summary_note, font_size=18, color=GRAY, italic=True,
        )

    # Progress dots at bottom
    dots = ""
    for i in range(total_sections):
        dots += "\u25CF  " if i <= section_idx else "\u25CB  "
    _add_textbox(
        slide,
        Inches(1), Inches(6.6), Inches(11.333), Inches(0.4),
        dots.strip(), font_size=10, color=LIGHT_GRAY,
    )


def _add_image_to_slide(slide, image_bytes: bytes, left, top, max_width, max_height):
    """Add an image to a slide, scaling to fit within max dimensions while keeping aspect ratio."""
    from PIL import Image

    img = Image.open(io.BytesIO(image_bytes))
    img_w, img_h = img.size

    # Calculate scale to fit within bounds
    scale_w = max_width / img_w
    scale_h = max_height / img_h
    scale = min(scale_w, scale_h)

    width = Emu(int(img_w * scale * 914400))
    height = Emu(int(img_h * scale * 914400))

    if width > max_width:
        width = max_width
    if height > max_height:
        height = max_height

    image_stream = io.BytesIO(image_bytes)
    slide.shapes.add_picture(image_stream, left, top, width, height)


def _add_content_slide(prs, section_heading: str, sub_heading: str, annotation: str,
                       bullets: list[BulletOut], image_data: bytes | None = None,
                       section_idx: int = 0, total_sections: int = 1):
    """Content slide with visual hierarchy, section badge, and typed bullet styling."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    has_image = image_data is not None
    text_width = Inches(7) if has_image else Inches(11)

    # Top accent line
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.04), PRIMARY)

    # Section badge with progress
    _add_section_badge(slide, section_heading, section_idx, total_sections)

    # Subsection heading
    _add_textbox(
        slide,
        Inches(0.8), Inches(1.0), text_width, Inches(0.7),
        sub_heading, font_size=26, color=HEADING, bold=True,
    )

    # Thin rule under heading
    _add_thin_rule(slide, Inches(0.8), Inches(1.75), Inches(2.5))

    # Annotation (speaker note preview)
    y_offset = Inches(1.95)
    if annotation:
        _add_textbox(
            slide,
            Inches(0.8), y_offset, text_width, Inches(0.7),
            annotation, font_size=13, color=GRAY, italic=True,
        )
        y_offset = Inches(2.7)

    # Bullets with type-aware styling
    if bullets:
        txBox = slide.shapes.add_textbox(Inches(0.8), y_offset, text_width, Inches(4.3))
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            # Determine bullet style based on content
            is_formula = bullet.latex is not None
            if is_formula:
                prefix = f"{FORMULA_CHAR}  "
                text_color = FORMULA_COLOR
            else:
                prefix = f"{BULLET_CHAR}  "
                text_color = DARK

            p.text = f"{prefix}{bullet.text}"
            p.font.size = Pt(16)
            p.font.color.rgb = text_color
            p.font.name = "Calibri"
            p.space_after = Pt(10)
            p.space_before = Pt(2)

            if is_formula:
                p.font.italic = True

    # Image in right column
    if has_image:
        try:
            _add_image_to_slide(
                slide,
                image_data,
                left=Inches(8.2),
                top=Inches(1.9),
                max_width=Inches(4.5),
                max_height=Inches(4.8),
            )
        except Exception:  # noqa: BLE001
            pass

    return slide


def generate_pptx(deck: DeckOut, image_loader=None) -> bytes:
    """Generate a PPTX file from a DeckOut structure. Returns bytes.

    image_loader: optional callable(storage_key) -> bytes that loads image data.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    total_sections = len(deck.sections)

    # Count total content slides for numbering
    total_slides = 1  # title
    for section in deck.sections:
        total_slides += 1  # section divider
        for sub in section.subsections:
            n_bullets = max(1, len(sub.bullets))
            total_slides += (n_bullets + MAX_BULLETS_PER_SLIDE - 1) // MAX_BULLETS_PER_SLIDE

    # Title slide
    _add_title_slide(prs, deck.title)

    slide_num = 1
    for s_idx, section in enumerate(deck.sections):
        # Section divider
        _add_section_slide(prs, section.heading, section.summary_note, s_idx, total_sections)
        slide_num += 1

        for sub in section.subsections:
            # Collect bullet objects (not just text) for type-aware rendering
            subsection_image: bytes | None = None

            if image_loader:
                for b in sub.bullets:
                    if b.image_url:
                        try:
                            subsection_image = image_loader(b.image_url)
                        except Exception:  # noqa: BLE001
                            pass
                        break

            # Split into multiple slides if too many bullets
            for page_start in range(0, max(1, len(sub.bullets)), MAX_BULLETS_PER_SLIDE):
                page_bullets = sub.bullets[page_start : page_start + MAX_BULLETS_PER_SLIDE]
                ann = sub.annotation if page_start == 0 else ""
                img = subsection_image if page_start == 0 else None

                content_slide = _add_content_slide(
                    prs, section.heading, sub.heading, ann,
                    page_bullets, image_data=img,
                    section_idx=s_idx, total_sections=total_sections,
                )
                slide_num += 1
                _add_slide_number(content_slide, slide_num, total_slides)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
